#!/usr/bin/env python3
import os
import json
import base64
import codecs
import hashlib
import datetime
import re
import shutil
import subprocess
import tempfile
import time
import anthropic

INBOX_MD      = "meta/inbox.md"
LOG_MD        = "meta/log.md"
INDEX_MD      = "meta/index.md"
ENTITIES_MD   = "meta/entities.md"
RAW_INBOX     = "raw/inbox"
RAW_PROCESSED = "raw/processed"

VALID_ZONE_PREFIXES = (
    "01-standards/",
    "02-workstreams/",
    "03-intelligence/",
    "04-internal/",
    "meta/",
)

# Pipeline state files the ingest agent must never rewrite. meta/ is otherwise
# writable (entities.md and index.md are legitimately model-maintained), but
# these two are written by this script alone: a model-generated "update" to
# inbox.md would wipe the dedup table and cause every processed file to be
# re-ingested, and one to log.md would erase the ingest history.
PROTECTED_PATHS = frozenset({"meta/inbox.md", "meta/log.md"})

# Maps a page's path prefix to the index heading it belongs under.
# Order matters: it doubles as the priority order when a wikilink target
# matches pages in more than one folder (e.g. Workstreams/OTC.md vs
# Open-Questions/OTC.md).
INDEX_SECTIONS = (
    ("01-standards/coding/",           "## 01-standards",    "### Coding"),
    ("01-standards/architecture/",     "## 01-standards",    "### Architecture"),
    ("01-standards/landscape/",        "## 01-standards",    "### Landscape"),
    ("02-workstreams/Workstreams/",    "## 02-workstreams",  "### Workstreams"),
    ("02-workstreams/Stakeholders/",   "## 02-workstreams",  "### Stakeholders"),
    ("02-workstreams/Meetings/",       "## 02-workstreams",  "### Meetings"),
    ("02-workstreams/Decisions/",      "## 02-workstreams",  "### Recent Decisions"),
    ("02-workstreams/Specs/",          "## 02-workstreams",  "### Specs"),
    ("02-workstreams/Developments/",   "## 02-workstreams",  "### Developments"),
    ("02-workstreams/Estimations/",    "## 02-workstreams",  "### Estimations"),
    ("02-workstreams/Issues/",         "## 02-workstreams",  "### Issues"),
    ("02-workstreams/Open-Questions/", "## 02-workstreams",  "### Open Questions"),
    ("03-intelligence/patterns/",        "## 03-intelligence", "### Patterns"),
    ("03-intelligence/lessons-learned/", "## 03-intelligence", "### Lessons Learned"),
    ("03-intelligence/gotchas/",         "## 03-intelligence", "### Gotchas"),
    ("03-intelligence/troubleshooting/", "## 03-intelligence", "### Troubleshooting Guides"),
    ("03-intelligence/faqs/",            "## 03-intelligence", "### FAQs"),
    ("04-internal/",                   "## 04-internal",     None),
)

ANTHROPIC_KEY = os.environ.get("ANTHROPIC_API_KEY")
if not ANTHROPIC_KEY:
    print("ERROR: ANTHROPIC_API_KEY not set")
    exit(1)

SINGLE_PASS_LIMIT = 150_000
CHUNK_SIZE        = 120_000
CHUNK_OVERLAP     =  10_000

MODEL             = "claude-opus-4-8"
MAX_OUTPUT_TOKENS = 64_000
API_RETRIES       = 3

# Prompt context budgets (characters)
PAGES_FULL_BUDGET   = 150_000  # full page content shown to Claude (updates need this)
PAGE_PREVIEW_CHARS  = 500      # fallback preview once the full budget is spent
INDEX_BUDGET        = 10_000
ENTITIES_BUDGET     = 16_000

def git_pull():
    try:
        result = subprocess.run(
            ["git", "pull", "--rebase"],
            capture_output=True, text=True, timeout=30
        )
        print(f"  git pull: {result.stdout.strip() or 'ok'}")
    except Exception as e:
        print(f"  git pull warning: {e}")

def get_payload():
    event_path = os.environ.get("GITHUB_EVENT_PATH")
    if not event_path:
        return None
    try:
        with open(event_path) as f:
            event = json.load(f)
        payload = event.get("client_payload", {})
        if payload.get("filename"):
            return payload
    except Exception as e:
        print(f"Could not read event payload: {e}")
    return None

# Image formats Claude vision accepts, mapped to their API media type.
# bmp/tiff are not supported by the API — they fall through to "unprocessable".
IMAGE_MEDIA_TYPES = {
    "png":  "image/png",
    "jpg":  "image/jpeg",
    "jpeg": "image/jpeg",
    "gif":  "image/gif",
    "webp": "image/webp",
}
IMAGE_MAX_BYTES = 4_500_000  # API limit is 5MB per image; leave headroom
AUDIO_VIDEO_TYPES = {
    "mp3", "wav", "m4a", "aac", "ogg", "flac", "wma", "opus",
    "mp4", "mov", "avi", "mkv", "webm", "wmv", "m4v", "flv",
}
# Legacy binary Office formats python-docx/python-pptx cannot read.
# If LibreOffice (soffice) is available they are converted to the modern
# format first; otherwise they are left for the curator with a clear message.
# (.xls is not in this set: xlrd reads it natively — see extract_xls.)
LEGACY_OFFICE_TYPES = {"doc", "ppt"}
LEGACY_TO_MODERN = {"doc": "docx", "ppt": "pptx", "xls": "xlsx"}
BINARY_EXTRACTORS = {
    "pdf":  "extract_pdf",
    "pptx": "extract_pptx",
    "docx": "extract_docx",
    "xlsx": "extract_xlsx",
    "xls":  "extract_xls",
}

def get_file_hash(filepath):
    try:
        with open(filepath, "rb") as f:
            return hashlib.md5(f.read()).hexdigest()
    except Exception:
        return None

def inbox_safe(filename):
    # "|" would corrupt the markdown table in meta/inbox.md
    return filename.replace("|", "¦")

def get_processed_records():
    records = {}
    try:
        with open(INBOX_MD, "r") as f:
            for line in f:
                parts = [p.strip() for p in line.split("|")]
                if len(parts) >= 4 and parts[1] and "." in parts[1]:
                    filename = parts[1].strip()
                    file_hash = parts[3].strip() if len(parts) > 3 else ""
                    records[filename] = file_hash
    except FileNotFoundError:
        pass
    return records

def is_already_processed(filename, filepath):
    records = get_processed_records()
    filename = inbox_safe(filename)
    if filename not in records:
        return False
    stored_hash = records[filename]
    current_hash = get_file_hash(filepath)
    if stored_hash and current_hash and stored_hash == current_hash:
        print(f"  Already processed (same content) — skipping.")
        return True
    print(f"  Same filename but different content — reprocessing.")
    return False

def scan_inbox_for_unprocessed():
    if not os.path.isdir(RAW_INBOX):
        return []
    records = get_processed_records()
    result = []
    for f in os.listdir(RAW_INBOX):
        if f.startswith(".") or f == "README.md":
            continue
        filepath = os.path.join(RAW_INBOX, f)
        key = inbox_safe(f)
        if key not in records:
            result.append(f)
        else:
            stored_hash = records[key]
            current_hash = get_file_hash(filepath)
            if stored_hash and current_hash and stored_hash != current_hash:
                result.append(f)
    return result

def strip_vtt(content):
    """Strip WEBVTT headers, cue numbers, timestamps, and NOTE blocks."""
    lines = []
    in_note = False
    for line in content.splitlines():
        stripped = line.strip()
        if stripped.startswith(("WEBVTT", "NOTE", "STYLE", "REGION")):
            in_note = stripped.startswith(("NOTE", "STYLE", "REGION"))
            continue
        if not stripped:
            in_note = False
            continue
        if in_note:
            continue
        if "-->" in stripped:
            continue
        if stripped.isdigit():
            continue
        # drop inline cue tags like <v Speaker Name> but keep the speaker
        stripped = re.sub(r"<v\s+([^>]+)>", r"\1: ", stripped)
        stripped = re.sub(r"</?[^>]+>", "", stripped)
        if stripped and (not lines or lines[-1] != stripped):
            lines.append(stripped)
    return "\n".join(lines)

def read_inbox_file(filename):
    path = os.path.join(RAW_INBOX, filename)
    if not os.path.exists(path):
        print(f"  ⚠ File not found: {path}")
        return None

    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    if ext in BINARY_EXTRACTORS:
        extractor = globals()[BINARY_EXTRACTORS[ext]]
        return extractor(path, filename)

    content = read_as_text(path, filename)
    if content and ext == "vtt":
        content = strip_vtt(content)
        print(f"  ✓ Stripped VTT timestamps — {len(content)} chars of transcript text")
    return content

# Byte-order marks, longest first: the UTF-32-LE BOM starts with the two
# bytes that are themselves the complete UTF-16-LE BOM, so checking UTF-16
# first would mis-decode every UTF-32-LE file.
BOM_ENCODINGS = (
    (codecs.BOM_UTF32_LE, "utf-32"),
    (codecs.BOM_UTF32_BE, "utf-32"),
    (codecs.BOM_UTF8,     "utf-8-sig"),
    (codecs.BOM_UTF16_LE, "utf-16"),
    (codecs.BOM_UTF16_BE, "utf-16"),
)
BINARY_SNIFF_BYTES = 8192
BINARY_CTRL_RATIO  = 0.30

def looks_binary(raw):
    """True if these bytes are not text in any single-byte or UTF-8 encoding.

    Callers must resolve BOM-declared UTF-16/32 first — those legitimately
    contain NUL bytes and would be misjudged here.
    """
    sample = raw[:BINARY_SNIFF_BYTES]
    if not sample:
        return False
    if b"\x00" in sample:
        return True
    # Control bytes outside tab/newline/carriage-return/formfeed/escape.
    ctrl = sum(1 for b in sample if b < 0x09 or 0x0E <= b < 0x20 or b == 0x7F)
    return ctrl / len(sample) > BINARY_CTRL_RATIO

def decode_text(raw):
    """Decode raw bytes to text. Returns (text, encoding_label), or
    (None, None) when the bytes are binary rather than text.

    Without the binary check, latin-1 with errors="replace" always succeeds —
    so an unrecognised binary format (.zip, .msg, .vsd) would decode to
    mojibake and be sent to the API as if it were a document.
    """
    for bom, enc in BOM_ENCODINGS:
        if raw.startswith(bom):
            try:
                return raw.decode(enc), enc
            except UnicodeDecodeError:
                break
    if looks_binary(raw):
        return None, None
    try:
        return raw.decode("utf-8"), "utf-8"
    except UnicodeDecodeError:
        return raw.decode("latin-1", errors="replace"), "latin-1"

def read_as_text(path, filename):
    try:
        with open(path, "rb") as f:
            raw = f.read()
    except Exception as e:
        print(f"  ⚠ Could not read {filename}: {e}")
        return None

    content, encoding = decode_text(raw)
    if content is None:
        print(f"  ⚠ {filename} is a binary format, not text — not sending to the API")
        return None
    print(f"  ✓ Read {len(content)} chars as {encoding}")
    return content

def extract_pdf(path, filename):
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            page_count = len(pdf.pages)
            pages = []
            for i, page in enumerate(pdf.pages):
                page_parts = []
                # page text first (it carries the document reading order,
                # including table words inline), structured [TABLE] blocks after
                text = page.extract_text()
                if text and text.strip():
                    page_parts.append(text.strip())
                try:
                    tables = page.extract_tables()
                    for table in tables:
                        if not table:
                            continue
                        rows = []
                        for row in table:
                            # keep empty cells as placeholders so values stay
                            # in their columns; only trim trailing empties
                            cells = ["" if cell is None else str(cell).strip()
                                     for cell in row]
                            while cells and not cells[-1]:
                                cells.pop()
                            if any(cells):
                                rows.append(" | ".join(cells))
                        if rows:
                            page_parts.append("[TABLE]\n" + "\n".join(rows))
                except Exception:
                    pass
                if page_parts:
                    pages.append(f"[Page {i+1}]\n" + "\n\n".join(page_parts))
            content = "\n\n".join(pages)
        if content:
            print(f"  ✓ Extracted {page_count} PDF pages, {len(content)} chars")
            return content
        print(f"  ⚠ PDF has no text layer ({page_count} pages) — attempting OCR")
        return extract_pdf_ocr(path)
    except ImportError:
        print("  ⚠ pdfplumber not installed — cannot extract PDF")
        return None
    except Exception as e:
        print(f"  ⚠ PDF extraction failed: {e}")
        return None

def extract_pdf_ocr(path):
    """Fallback for scanned PDFs with no text layer: rasterize each page and
    OCR it. Needs pytesseract + pdf2image (and the tesseract/poppler binaries);
    without them the file is left for the curator as before."""
    try:
        import pytesseract
        from pdf2image import convert_from_path
    except ImportError:
        print("  ⚠ OCR unavailable (pytesseract/pdf2image not installed) — needs manual handling")
        return None
    try:
        images = convert_from_path(path, dpi=200)
        pages = []
        for i, image in enumerate(images):
            text = pytesseract.image_to_string(image)
            if text.strip():
                pages.append(f"[Page {i+1} (OCR)]\n{text.strip()}")
        content = "\n\n".join(pages)
        if content:
            print(f"  ✓ OCR extracted {len(pages)} of {len(images)} pages, {len(content)} chars")
            return content
        print(f"  ⚠ OCR found no text in {len(images)} pages")
        return None
    except Exception as e:
        print(f"  ⚠ OCR failed: {e}")
        return None

def pptx_shape_texts(shape, texts):
    # grouped shapes hold their content in child shapes — recurse into them
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            pptx_shape_texts(child, texts)
        return
    if shape.has_table:
        rows = []
        for row in shape.table.rows:
            # keep empty cells as placeholders so values stay in their
            # columns; only trim trailing empties
            cells = [cell.text.strip() for cell in row.cells]
            while cells and not cells[-1]:
                cells.pop()
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            texts.append("[TABLE]\n" + "\n".join(rows))
    elif hasattr(shape, "text_frame") and shape.text_frame:
        text = shape.text_frame.text.strip()
        if text:
            texts.append(text)
    elif hasattr(shape, "text") and shape.text.strip():
        texts.append(shape.text.strip())

def extract_pptx_markitdown(path):
    """Preferred PPTX path: markitdown emits one markdown block per slide
    (with '<!-- Slide number: N -->' markers), handles grouped shapes,
    renders tables as markdown with empty cells preserved, and includes
    speaker notes under '### Notes:'."""
    try:
        from markitdown import MarkItDown
    except ImportError:
        return None
    try:
        text = MarkItDown().convert(path).text_content
        if text and text.strip():
            return text
    except Exception as e:
        print(f"  ⚠ markitdown failed ({e}) — falling back to python-pptx")
    return None

def extract_pptx(path, filename):
    content = extract_pptx_markitdown(path)
    if content:
        print(f"  ✓ Extracted {len(content)} chars from PPTX via markitdown")
        return content
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                pptx_shape_texts(shape, texts)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    texts.append(f"[Speaker notes]\n{notes}")
            if texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        content = "\n\n".join(slides)
        if content:
            print(f"  ✓ Extracted {len(prs.slides)} PPTX slides, {len(content)} chars")
            return content
        print(f"  ⚠ PPTX has no extractable text ({len(prs.slides)} slides)")
        return None
    except ImportError:
        print("  ⚠ python-pptx not installed — cannot extract PPTX")
        return None
    except Exception as e:
        print(f"  ⚠ PPTX extraction failed: {e}")
        return None

def docx_table_rows(table):
    rows = []
    for row in table.rows:
        cells = []
        prev_tc = None
        for cell in row.cells:
            # horizontally merged cells repeat the same underlying element —
            # emit the text once, not once per spanned column
            if cell._tc is prev_tc:
                continue
            prev_tc = cell._tc
            cells.append(cell.text.strip())
        # keep empty cells as placeholders so values stay in their columns;
        # only trim trailing empties
        while cells and not cells[-1]:
            cells.pop()
        if any(cells):
            rows.append(" | ".join(cells))
    return rows

def extract_docx_pandoc(path):
    """Preferred DOCX path: pandoc produces full-fidelity markdown — heading
    levels, lists, links, emphasis, and tables that keep empty cells and
    merged-cell spans (colspan) intact, all in document order."""
    if not shutil.which("pandoc"):
        return None
    try:
        result = subprocess.run(
            ["pandoc", "-t", "gfm", "--wrap=none", path],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
        if result.returncode != 0:
            print(f"  ⚠ pandoc failed ({result.stderr.strip()[:200]}) — falling back to python-docx")
    except Exception as e:
        print(f"  ⚠ pandoc error ({e}) — falling back to python-docx")
    return None

def extract_docx(path, filename):
    content = extract_docx_pandoc(path)
    if content:
        print(f"  ✓ Extracted {len(content)} chars from DOCX via pandoc")
        return content
    try:
        from docx import Document
        from docx.table import Table
        from docx.text.paragraph import Paragraph
        doc = Document(path)
        parts = []
        # walk the document body in order so tables stay next to the
        # headings and paragraphs they belong to
        for child in doc.element.body.iterchildren():
            if child.tag.endswith("}p"):
                para = Paragraph(child, doc)
                text = para.text.strip()
                if not text:
                    continue
                style = para.style.name if para.style else ""
                if style.startswith("Heading"):
                    try:
                        level = min(int(style.rsplit(" ", 1)[-1]), 6)
                    except ValueError:
                        level = 1
                    text = "#" * level + " " + text
                parts.append(text)
            elif child.tag.endswith("}tbl"):
                rows = docx_table_rows(Table(child, doc))
                if rows:
                    parts.append("[TABLE]\n" + "\n".join(rows))
        content = "\n".join(parts)
        if content:
            print(f"  ✓ Extracted {len(content)} chars from DOCX")
            return content
        print("  ⚠ DOCX has no extractable text")
        return None
    except ImportError:
        print("  ⚠ python-docx not installed — cannot extract DOCX")
        return None
    except Exception as e:
        print(f"  ⚠ DOCX extraction failed: {e}")
        return None

def extract_xlsx(path, filename):
    try:
        import openpyxl
        # data_only=True gives cached formula results; the raw load lets us
        # fall back to the formula text when no cached value exists (e.g.
        # workbooks generated by tools other than Excel)
        wb = openpyxl.load_workbook(path, data_only=True)
        wb_raw = openpyxl.load_workbook(path, data_only=False)
        sheets = []
        uncached_formulas = 0
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            ws_raw = wb_raw[sheet_name]
            rows = []
            for row, row_raw in zip(ws.iter_rows(values_only=True),
                                    ws_raw.iter_rows(values_only=True)):
                cells = []
                for c, c_raw in zip(row, row_raw):
                    if c is None and isinstance(c_raw, str) and c_raw.startswith("="):
                        cells.append(c_raw)
                        uncached_formulas += 1
                    elif c is None:
                        # keep empty cells as placeholders so values stay
                        # in their columns
                        cells.append("")
                    else:
                        cells.append(str(c))
                while cells and not cells[-1].strip():
                    cells.pop()
                if any(c.strip() for c in cells):
                    rows.append(" | ".join(cells))
            if rows:
                sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        if uncached_formulas:
            print(f"  ⚠ {uncached_formulas} formula cell(s) had no cached value — emitted the formula text instead")
        content = "\n\n".join(sheets)
        if content:
            print(f"  ✓ Extracted {len(wb.sheetnames)} XLSX sheets, {len(content)} chars")
            return content
        print("  ⚠ XLSX has no extractable data")
        return None
    except ImportError:
        print("  ⚠ openpyxl not installed — cannot extract XLSX")
        return None
    except Exception as e:
        print(f"  ⚠ XLSX extraction failed: {e}")
        return None

def extract_xls(path, filename):
    """Legacy .xls: read natively with xlrd — no LibreOffice needed. The BIFF
    format stores cached formula results, so formula cells yield their values.
    Falls back to LibreOffice conversion for files xlrd cannot parse."""
    try:
        import xlrd
    except ImportError:
        print("  ⚠ xlrd not installed — trying LibreOffice conversion")
        return extract_legacy_office(path, filename, "xls")
    try:
        wb = xlrd.open_workbook(path)
        sheets = []
        for ws in wb.sheets():
            rows = []
            for r in range(ws.nrows):
                cells = []
                for c in range(ws.ncols):
                    cell = ws.cell(r, c)
                    if cell.ctype in (xlrd.XL_CELL_EMPTY, xlrd.XL_CELL_BLANK):
                        # keep empty cells as placeholders so values stay
                        # in their columns
                        cells.append("")
                    elif (cell.ctype == xlrd.XL_CELL_NUMBER
                          and cell.value == int(cell.value)):
                        cells.append(str(int(cell.value)))
                    else:
                        cells.append(str(cell.value).strip())
                while cells and not cells[-1]:
                    cells.pop()
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                sheets.append(f"[Sheet: {ws.name}]\n" + "\n".join(rows))
        content = "\n\n".join(sheets)
        if content:
            print(f"  ✓ Extracted {wb.nsheets} XLS sheets via xlrd, {len(content)} chars")
            return content
        print("  ⚠ XLS has no extractable data via xlrd — trying LibreOffice conversion")
    except Exception as e:
        print(f"  ⚠ xlrd failed ({e}) — trying LibreOffice conversion")
    return extract_legacy_office(path, filename, "xls")

def extract_legacy_office(path, filename, ext):
    """Convert a legacy .doc/.ppt/.xls to its modern format with LibreOffice,
    then run the normal extractor on the result. Returns None when soffice is
    not installed or the conversion fails."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None
    target = LEGACY_TO_MODERN[ext]
    outdir = tempfile.mkdtemp(prefix="vault-legacy-")
    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", target, "--outdir", outdir, path],
            capture_output=True, text=True, timeout=180,
        )
        base = os.path.splitext(os.path.basename(path))[0]
        converted = os.path.join(outdir, base + "." + target)
        if result.returncode != 0 or not os.path.exists(converted):
            print(f"  ⚠ LibreOffice conversion failed: {(result.stderr or result.stdout).strip()[:200]}")
            return None
        print(f"  ✓ Converted legacy .{ext} to .{target} via LibreOffice")
        extractor = globals()[BINARY_EXTRACTORS[target]]
        return extractor(converted, filename)
    except Exception as e:
        print(f"  ⚠ LibreOffice conversion error: {e}")
        return None
    finally:
        shutil.rmtree(outdir, ignore_errors=True)

def move_to_processed(filename):
    os.makedirs(RAW_PROCESSED, exist_ok=True)
    src = os.path.join(RAW_INBOX, filename)
    dst = os.path.join(RAW_PROCESSED, filename)
    if os.path.exists(src):
        shutil.move(src, dst)
        print(f"  ✓ Moved to processed: {filename}")
    else:
        print(f"  ⚠ Could not move — file not found: {src}")

def is_valid_vault_path(path):
    if not path or not path.strip():
        return False
    path = path.replace("\\", "/").strip()
    if os.path.isabs(path):
        return False
    normalized = os.path.normpath(path).replace("\\", "/")
    if normalized.startswith("..") or "/../" in normalized:
        return False
    if normalized in PROTECTED_PATHS:
        return False
    for prefix in VALID_ZONE_PREFIXES:
        if normalized.startswith(prefix):
            return True
    return False

def build_prompt_context():
    try:
        with open("CLAUDE.md", "r") as f:
            schema = f.read()
    except FileNotFoundError:
        schema = "Follow standard ABAP vault conventions."

    try:
        with open(INDEX_MD, "r") as f:
            index = f.read()
    except FileNotFoundError:
        index = ""

    try:
        with open(ENTITIES_MD, "r") as f:
            entities = f.read()
    except FileNotFoundError:
        entities = ""

    # Every content page in the vault, with FULL content up to the budget.
    # Claude may only UPDATE a page it has seen in full — an update from a
    # truncated preview would silently destroy the unseen part of the page.
    page_paths = []
    for zone_dir in ["01-standards", "02-workstreams", "03-intelligence", "04-internal"]:
        if os.path.isdir(zone_dir):
            for root, dirs, files in os.walk(zone_dir):
                for f in sorted(files):
                    if f.endswith(".md") and not f.startswith("_") and f != "README.md":
                        page_paths.append(os.path.join(root, f).replace("\\", "/"))

    existing_pages = "## Complete page listing\n" + "\n".join(f"- {p}" for p in page_paths) + "\n"
    remaining = PAGES_FULL_BUDGET
    previews = []
    for p in page_paths:
        try:
            with open(p, "r") as fh:
                content = fh.read()
        except Exception:
            continue
        if len(content) <= remaining:
            existing_pages += f"\n### FULL: {p}\n{content}\n"
            remaining -= len(content)
        else:
            previews.append(f"\n### PREVIEW ONLY (do not update): {p}\n{content[:PAGE_PREVIEW_CHARS]}\n")
    if previews:
        print(f"  ⚠ Page-context budget exhausted — {len(previews)} page(s) sent as preview only")
        existing_pages += "".join(previews)

    return schema, index, entities, existing_pages

def extract_response_text(response):
    for block in response.content:
        if block.type == "text":
            return block.text
    return ""

def call_api_with_retry(client, messages):
    """One streamed API call with backoff on retryable errors.
    Returns the final Message, or None after exhausting retries / on truncation."""
    last_err = None
    for attempt in range(API_RETRIES):
        if attempt:
            wait = 15 * (2 ** attempt)
            print(f"  ⚠ API error ({last_err}) — retrying in {wait}s ({attempt + 1}/{API_RETRIES})")
            time.sleep(wait)
        try:
            with client.messages.stream(
                model=MODEL,
                max_tokens=MAX_OUTPUT_TOKENS,
                thinking={"type": "adaptive"},
                messages=messages,
            ) as stream:
                response = stream.get_final_message()
            if response.stop_reason == "max_tokens":
                print(f"  ⚠ Response truncated at {MAX_OUTPUT_TOKENS} output tokens — treating as failure")
                return None
            return response
        except anthropic.RateLimitError as e:
            last_err = e
        except anthropic.APIStatusError as e:
            if e.status_code < 500:
                print(f"  ⚠ Non-retryable API error {e.status_code}: {e.message}")
                return None
            last_err = e
        except anthropic.APIConnectionError as e:
            last_err = e
    print(f"  ⚠ API failed after {API_RETRIES} attempts: {last_err}")
    return None

def parse_result_json(text):
    """Extract and parse the JSON object from a response. Returns (result, error)."""
    code_match = re.search(r'```(?:json)?\s*(\{[\s\S]*?\})\s*```', text)
    if code_match:
        json_str = code_match.group(1)
    else:
        json_match = re.search(r'\{[\s\S]*\}', text)
        if not json_match:
            return None, "no JSON object found in the response"
        json_str = json_match.group()
    try:
        return json.loads(json_str), None
    except json.JSONDecodeError as e:
        return None, str(e)

def call_claude(filename, content_chunk, schema, index, entities, existing_pages,
                chunk_info="", is_continuation=False, image_path=None, image_media_type=None):
    client = anthropic.Anthropic(api_key=ANTHROPIC_KEY)

    continuation_note = ""
    if is_continuation:
        continuation_note = (
            "NOTE: This is a continuation chunk. Previous chunk already processed. "
            "Focus on NEW information. UPDATE existing pages if needed. "
            "Do not re-create pages from earlier chunks.\n\n"
        )

    today = datetime.date.today().isoformat()

    prompt = f"""You are the ABAP Vault AI, the disciplined librarian maintaining the ABAP project's knowledge vault.

## Vault Constitution (CLAUDE.md)
{schema}

## Current Index
{index[:INDEX_BUDGET]}

## Entity Registry (canonical names — meta/entities.md)
{entities[:ENTITIES_BUDGET]}

## Existing Pages (do not duplicate these)
{existing_pages}

## File to Ingest
Filename: `{filename}`{chunk_info}

{continuation_note}Content:
---
{content_chunk}
---

CRITICAL RULES — follow these exactly or the vault breaks:

1. ZONE PATHS: Every page path MUST start with its zone folder. Valid path prefixes ONLY:
   - 01-standards/coding/
   - 01-standards/architecture/
   - 01-standards/landscape/
   - 02-workstreams/Workstreams/
   - 02-workstreams/Stakeholders/{{WS}}/
   - 02-workstreams/Meetings/{{WS}}/
   - 02-workstreams/Decisions/{{WS}}/
   - 02-workstreams/Specs/{{WS}}/
   - 02-workstreams/Developments/{{WS}}/
   - 02-workstreams/Estimations/{{WS}}/
   - 02-workstreams/Issues/{{WS}}/
   - 02-workstreams/Open-Questions/
   - 03-intelligence/patterns/
   - 03-intelligence/lessons-learned/
   - 03-intelligence/gotchas/
   - 03-intelligence/troubleshooting/
   - 03-intelligence/faqs/technical/
   - 03-intelligence/faqs/process-and-transport/
   - 03-intelligence/faqs/landscape-and-access/
   - 03-intelligence/faqs/testing/
   - 04-internal/contacts/
   - 04-internal/onboarding/
   - 04-internal/processes/
   - 04-internal/runbooks/
   - meta/
   NEVER write to the vault root. NEVER invent new folders.

2. ENTITY NORMALIZATION (mandatory before creating anything workstream-related):
   {{WS}} is ALWAYS the canonical workstream slug from the Entity Registry above.
   Normalize every workstream, module, system, and vendor name found in the document
   (lowercase, strip punctuation, resolve abbreviations) and match it against the
   registry's aliases. Exact or fuzzy match → use the canonical slug. No match →
   add the new entity to meta/entities.md FIRST (as an update), then create pages.
   When uncertain, treat as an existing entity. Never invent a second spelling.
   MODULE-SLUG RULE: for legacy objects not tied to any active workstream, {{WS}}
   may instead be a canonical MODULE slug from the registry (SD, MM, FI, ...),
   e.g. 02-workstreams/Developments/SD/SD - ZSD_ORDER_CHECK.md.

3. FRONTMATTER: Every page MUST start with exactly this YAML frontmatter:
   ---
   title: ""
   type: ""
   zone: ""
   status: active
   owner: ""
   created: {today}
   updated: {today}
   workstream: ""
   tags: []
   source_files: []
   ---
   Valid type values: standard, architecture, landscape, workstream, stakeholder,
   meeting, decision, spec, development, estimation, issue, open-questions,
   pattern, lessons-learned, gotcha, troubleshooting, faq,
   contact, onboarding, process, runbook
   Valid zone values: 01-standards, 02-workstreams, 03-intelligence, 04-internal
   Valid status values: active, draft, parked, archived, resolved, evergreen
   Set workstream to the canonical slug for all Zone 02 pages. Add `{filename}` to
   source_files on every page this ingest touches.
   DO NOT invent new frontmatter fields. DO NOT omit required fields.
   TAG DISCIPLINE: every value in tags MUST come from the Tag Vocabulary section
   of the Entity Registry above. Tags are domain tags only (technology,
   business-object, quality, process, role, phase) — NEVER echo the page's type
   or workstream into tags. A genuinely new tag → add it to the Tag Vocabulary
   in meta/entities.md FIRST (as an update, with its category), then use it.
   Known alias (e.g. sm36, batch) → use the canonical tag (batch-job).

4. WIKILINKS: Every page MUST contain at least one [[wikilink]] to a related page,
   and every Zone 02 page MUST link to its workstream page [[{{WS}}]].

5. NO DUPLICATES: If a page already exists in Existing Pages above, UPDATE it.
   Do not create a new page for the same content. If a decision on the same topic
   exists, update it and note the revision — never create a competing decision page.

6. UPDATE ORDER (Zone 02 — non-negotiable): workstream page first, then stakeholders,
   then meeting/spec artifacts, then decisions, then developments, then issues,
   then the rolling Open-Questions/{{WS}}.md page, then promote to Zone 03.

7. MEETING CONVERSION RULE: Meetings are source material, not final artifacts.
   Only create a meeting page if it carries information beyond what you pushed into
   the parent pages, and list which pages it updated. Never create a page that is
   just a summary with no linked updates.

8. FAQ EXTRACTION RULE (mandatory for transcripts and email chains): For every
   question asked by the client, a functional consultant, or a team member —
   classify it into faqs/technical/, faqs/process-and-transport/,
   faqs/landscape-and-access/, or faqs/testing/. Answered in-session → append under
   ## Answered Questions with attribution. Unanswered → append to the
   ## Unanswered Questions table AND add it to 02-workstreams/Open-Questions/{{WS}}.md.
   Same question already recorded → add the new asker to `Also asked by`.

9. GOTCHA RULE: Non-obvious SAP/ABAP behavior that cost real time gets a
   03-intelligence/gotchas/Gotcha - {{Name}}.md page IMMEDIATELY — gotchas do not
   wait for a second occurrence.

10. PATTERN RULE: If the same approach, error, workaround, or tradeoff appears in
    2+ workstreams or 2+ developments, create or update
    03-intelligence/patterns/Pattern - {{Name}}.md linking every place it was observed.

11. FLOATING PAGES ARE FORBIDDEN: If you cannot link a new page upward to a parent,
    do not create it — append to an existing page instead.

12. CODE INGESTION RULE: If the input is ABAP source code (programs, classes,
    function modules, CDS views, enhancements), reconstruct the missing documentation:
    - CREATE/UPDATE a Spec page (Specs/{{WS}}/) with: functional summary in business
      terms, technical specification, and impacted business processes.
    - CREATE/UPDATE a Development record (Developments/{{WS}}/) with object details
      and a Dependencies section listing EVERY referenced object (tables, classes,
      function modules, CDS views, enhancements) as [[wikilinks]] — link them even
      if their pages don't exist yet.
    - All pages generated from code MUST have status: draft and tags: [ai-generated]
      — the business rationale is inferred, not confirmed. A human SME removes the tag.
    - Use the owning workstream slug if determinable, else the module slug (rule 2).
    - If the input is a dependency extract (TADIR, where-used list, CDS dependency
      export): update the Dependencies sections of the affected object pages.
      Do NOT create standalone dependency-dump pages.

13. TROUBLESHOOTING PROMOTION: If 2+ resolved Issues share the same diagnostic path,
    create or update 03-intelligence/troubleshooting/Troubleshooting - {{Area}}.md
    (symptoms, ordered diagnostic checks, known causes and fixes), linking every
    source Issue.

14. ESTIMATION RULE: Effort estimates go to Estimations/{{WS}}/{{WS}} - Estimation -
    {{Topic}} - {{YYYY-MM-DD}}.md with scope, assumptions, and a per-object table
    (estimated vs actual). When actuals appear in a document, update BOTH the
    estimation page and the effort fields of the affected Development records.

15. UPDATE SEMANTICS: When you UPDATE a page, return its COMPLETE content —
    every existing section plus your changes. Never drop or shorten sections
    you are not changing. Preserve the page's original `created:` date and set
    `updated: {today}`. Only update pages whose FULL content appears above
    (marked "FULL:"). Pages marked "PREVIEW ONLY" must NOT be updated — if
    one needs new information, note it in the log_entry instead.

16. NO DURABLE CONTENT: If the document contains no durable knowledge
    (logistics chatter, empty shells, generated code with no purpose), return
    empty "updates" and "creates" arrays with only a log_entry explaining why.
    Do not invent pages to feel useful.

Job:
1. Triage the document per the constitution (workstream-specific → Zone 02,
   reusable learning → Zone 03, standard/landscape → Zone 01, team ops → Zone 04)
2. UPDATE existing pages if new information exists
3. CREATE new pages following the naming rules in the constitution
4. Update meta/entities.md if new entities or aliases were found
5. Add index entries
6. Write log entry

Respond ONLY with valid JSON — no markdown fences, no explanation, just the JSON object:
{{
  "updates": [{{"path": "zone-folder/filename.md", "content": "full markdown with frontmatter"}}],
  "creates": [{{"path": "zone-folder/filename.md", "content": "full markdown with frontmatter"}}],
  "index_entries": ["- [[Page]] — description"],
  "log_entry": "- {today}: Ingested {filename}. Updated X. Created Y."
}}

Minimal example of a valid response (shape only — your content will differ):
{{
  "updates": [{{"path": "02-workstreams/Workstreams/OTC.md", "content": "---\\ntitle: \\"OTC\\"\\n...full page...\\n"}}],
  "creates": [],
  "index_entries": [],
  "log_entry": "- {today}: Ingested {filename}. Updated [[OTC]] with wave 2 scope."
}}"""

    if image_path:
        with open(image_path, "rb") as f:
            image_data = base64.standard_b64encode(f.read()).decode("utf-8")
        user_content = [
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": image_media_type,
                    "data": image_data,
                },
            },
            {"type": "text", "text": prompt},
        ]
    else:
        user_content = prompt

    messages = [{"role": "user", "content": user_content}]

    response = call_api_with_retry(client, messages)
    if response is None:
        return None

    text = extract_response_text(response)
    result, parse_err = parse_result_json(text)
    if result is not None:
        return result

    # One retry with the parse error fed back
    print(f"  ⚠ JSON parse error ({parse_err}) — retrying once with error feedback")
    retry_messages = messages + [
        {"role": "assistant", "content": text or "(empty response)"},
        {
            "role": "user",
            "content": (
                f"Your previous response could not be parsed as JSON: {parse_err}. "
                "Respond again with ONLY the complete, valid JSON object — "
                "no markdown fences, no commentary, properly escaped strings."
            ),
        },
    ]
    response = call_api_with_retry(client, retry_messages)
    if response is None:
        return None
    result, parse_err = parse_result_json(extract_response_text(response))
    if result is None:
        print(f"  ⚠ JSON parse failed again: {parse_err}")
    return result

def merge_results(results):
    merged = {"updates": [], "creates": [], "index_entries": [], "log_entry": ""}
    seen_updates = {}
    seen_creates = {}

    for result in results:
        if not result:
            continue
        for item in result.get("updates", []):
            seen_updates[item["path"]] = item
        for item in result.get("creates", []):
            seen_creates[item["path"]] = item
        for entry in result.get("index_entries", []):
            if entry not in merged["index_entries"]:
                merged["index_entries"].append(entry)
        if result.get("log_entry"):
            merged["log_entry"] = result["log_entry"]

    merged["updates"] = list(seen_updates.values())
    merged["creates"] = list(seen_creates.values())
    return merged

def run_ingest_agent(filename, content):
    schema, index, entities, existing_pages = build_prompt_context()

    if len(content) <= SINGLE_PASS_LIMIT:
        print(f"  Content: {len(content):,} chars — single pass")
        return call_claude(filename, content, schema, index, entities, existing_pages)

    chunks = []
    start = 0
    while start < len(content):
        end = start + CHUNK_SIZE
        chunks.append(content[start:end])
        if end >= len(content):
            break
        start = end - CHUNK_OVERLAP

    print(f"  Content: {len(content):,} chars — {len(chunks)} chunks")

    results = []
    for i, chunk in enumerate(chunks):
        chunk_info = f"\n(Chunk {i+1} of {len(chunks)})"
        print(f"  Processing chunk {i+1}/{len(chunks)}...")
        result = call_claude(
            filename, chunk, schema, index, entities, existing_pages,
            chunk_info=chunk_info,
            is_continuation=(i > 0)
        )
        if result:
            results.append(result)
            # log_only=True for EVERY chunk: pages are written so later chunks
            # see them, but the log/inbox entries are written exactly once —
            # by process_file applying the merged result.
            apply_vault_changes(result, filename, log_only=True)
            try:
                with open(INDEX_MD, "r") as f:
                    index = f.read()
            except Exception:
                pass
            _, _, entities, existing_pages = build_prompt_context()

    if len(results) <= 1:
        return results[0] if results else None

    return merge_results(results)

def wikilink_target(text):
    m = re.search(r"\[\[([^\]|]+)(?:\|[^\]]*)?\]\]", text)
    return m.group(1).strip() if m else None

def scan_vault_pages():
    pages = []
    for zone_dir in ("01-standards", "02-workstreams", "03-intelligence", "04-internal"):
        for root, dirs, files in os.walk(zone_dir):
            for f in files:
                if f.endswith(".md"):
                    pages.append(os.path.join(root, f).replace("\\", "/"))
    return pages

def find_page_path(target, touched_paths):
    if not target:
        return None
    candidates = [p.replace("\\", "/") for p in touched_paths] + scan_vault_pages()
    if "/" in target:
        for p in candidates:
            if p.endswith("/" + target + ".md") or p == target + ".md":
                return p
    matches = [p for p in candidates if os.path.basename(p) == target + ".md"]
    if not matches:
        return None
    for prefix, _, _ in INDEX_SECTIONS:
        for p in matches:
            if p.startswith(prefix):
                return p
    return matches[0]

def section_for_path(path):
    for prefix, zone_heading, section_heading in INDEX_SECTIONS:
        if path.startswith(prefix):
            return zone_heading, section_heading
    return None, None

def find_block(lines, heading, start=0, end=None):
    """Return (heading_index, block_end) for a heading, searching lines[start:end].
    block_end is the index of the next heading of same-or-higher level, or end."""
    if end is None:
        end = len(lines)
    level = heading.split(" ")[0]  # "##" or "###"
    for i in range(start, end):
        if lines[i].strip() == heading:
            for j in range(i + 1, end):
                stripped = lines[j].strip()
                if stripped.startswith("#") and len(stripped.split(" ")[0]) <= len(level):
                    return i, j
            return i, end
    return None, None

def insert_entry_in_section(lines, entry, zone_heading, section_heading):
    zone_start, zone_end = find_block(lines, zone_heading)
    if zone_start is None:
        lines.extend(["", zone_heading, ""])
        zone_start, zone_end = len(lines) - 2, len(lines)

    if section_heading:
        sec_start, sec_end = find_block(lines, section_heading, zone_start + 1, zone_end)
        if sec_start is None:
            insert_at = zone_end
            while insert_at > zone_start + 1 and not lines[insert_at - 1].strip():
                insert_at -= 1
            lines[insert_at:insert_at] = ["", section_heading, ""]
            sec_start, sec_end = insert_at + 1, insert_at + 3
    else:
        sec_start, sec_end = zone_start, zone_end

    # drop the "(none yet)" placeholder now that the section has a real entry
    for i in range(sec_end - 1, sec_start, -1):
        if lines[i].strip().startswith("- _(none"):
            del lines[i]
            sec_end -= 1

    insert_at = sec_start + 1
    for i in range(sec_start + 1, sec_end):
        if lines[i].strip().startswith("- "):
            insert_at = i + 1
    if insert_at == sec_start + 1 and insert_at < sec_end and not lines[insert_at].strip():
        insert_at += 1
    lines.insert(insert_at, entry)

def update_index(index_entries, touched_paths):
    try:
        with open(INDEX_MD, "r") as f:
            lines = f.read().splitlines()
    except FileNotFoundError:
        lines = ["# Vault Index", ""]

    today = datetime.date.today().isoformat()
    for i, line in enumerate(lines):
        if line.startswith("_Last updated:"):
            lines[i] = f"_Last updated: {today}_"
            break

    for entry in index_entries:
        entry = entry.strip()
        if not entry:
            continue
        target = wikilink_target(entry)

        # same page already indexed → refresh that line instead of duplicating
        replaced = False
        for i, line in enumerate(lines):
            if line.strip().startswith("- ") and target and wikilink_target(line) == target:
                lines[i] = entry
                replaced = True
                break
        if replaced:
            continue

        path = find_page_path(target, touched_paths)
        zone_heading, section_heading = section_for_path(path) if path else (None, None)
        if zone_heading is None:
            print(f"  ⚠ Index entry has no matching vault page, appending at end: {entry[:80]}")
            lines.extend(["", entry])
            continue
        insert_entry_in_section(lines, entry, zone_heading, section_heading)

    with open(INDEX_MD, "w") as f:
        f.write("\n".join(lines) + "\n")

def apply_vault_changes(result, filename, log_only=False):
    if not result:
        return False

    changed = False

    for item in result.get("updates", []) + result.get("creates", []):
        path = item.get("path", "")
        content = item.get("content", "")

        if not is_valid_vault_path(path):
            reason = ("pipeline state file, written by this script only"
                      if os.path.normpath(path).replace("\\", "/") in PROTECTED_PATHS
                      else "not a valid zone folder")
            print(f"  ⚠ Rejected path: '{path}' — {reason}")
            continue

        dir_name = os.path.dirname(path)
        if dir_name:
            os.makedirs(dir_name, exist_ok=True)
        with open(path, "w") as f:
            f.write(content)
        print(f"  ✓ Written: {path}")
        changed = True

    if result.get("index_entries"):
        try:
            touched = [i.get("path", "").replace("\\", "/") for i in
                       result.get("updates", []) + result.get("creates", [])]
            update_index(result["index_entries"], touched)
        except Exception as e:
            print(f"  ⚠ Index update error: {e}")

    if log_only:
        return changed

    if result.get("log_entry"):
        with open(LOG_MD, "a") as f:
            f.write(f"\n{result['log_entry']}\n")

    today = datetime.date.today().isoformat()
    file_hash = get_file_hash(os.path.join(RAW_INBOX, filename)) or "—"
    updated = ", ".join(i["path"] for i in result.get("updates", []))
    created = ", ".join(i["path"] for i in result.get("creates", []))
    with open(INBOX_MD, "a") as f:
        f.write(f"| {inbox_safe(filename)} | {today} | {file_hash} | {updated or '—'} | {created or '—'} |\n")

    return changed

def log_unprocessable(filename, reason, permanent=True):
    """Log a file we cannot process and leave it in the inbox for the curator.

    permanent=True  — the file itself is the problem (unsupported format, no
        extractable text). Recorded in meta/inbox.md against its current hash
        so the weekly sweep stops retrying it and stops re-appending this same
        line to the log every run. Replacing the file changes its hash, which
        makes it eligible again automatically.
    permanent=False — a transient failure (API error, agent gave up). NOT
        recorded, so the next push or sweep retries it as before.
    """
    scope = "left in inbox for the curator" if permanent else "will retry next run"
    print(f"  ⚠ Unprocessable — {scope}: {reason}")
    today = datetime.date.today().isoformat()
    with open(LOG_MD, "a") as f:
        f.write(f"\n- {today}: Could not ingest `{filename}` — {reason}. "
                f"File left in raw/inbox/ for the curator.\n")
    if permanent:
        file_hash = get_file_hash(os.path.join(RAW_INBOX, filename)) or "—"
        with open(INBOX_MD, "a") as f:
            f.write(f"| {inbox_safe(filename)} | {today} | {file_hash} | — | "
                    f"UNPROCESSABLE: {reason} |\n")

def process_image_file(filename, filepath, ext):
    if ext not in IMAGE_MEDIA_TYPES:
        log_unprocessable(filename, f"unsupported image format .{ext} — convert to PNG or JPG")
        return
    if os.path.getsize(filepath) > IMAGE_MAX_BYTES:
        log_unprocessable(filename, "image exceeds the 5MB API limit — resize and re-drop")
        return

    print(f"  Image — running ingest agent with vision...")
    schema, index, entities, existing_pages = build_prompt_context()
    placeholder = (
        f"[IMAGE FILE: {filename} — the image is attached above. It is likely an "
        "architecture diagram, whiteboard photo, or screenshot. Extract the durable "
        "knowledge it shows.]"
    )
    result = call_claude(
        filename, placeholder, schema, index, entities, existing_pages,
        image_path=filepath, image_media_type=IMAGE_MEDIA_TYPES[ext],
    )
    if not result:
        log_unprocessable(filename, "ingest agent failed on image (see run log)",
                          permanent=False)
        return
    apply_vault_changes(result, filename)
    move_to_processed(filename)
    print(f"  ✓ Done: {filename}")

def process_file(filename):
    print(f"\nProcessing: {filename}")

    filepath = os.path.join(RAW_INBOX, filename)

    if is_already_processed(filename, filepath):
        return

    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    if ext in IMAGE_MEDIA_TYPES or ext in ("bmp", "tiff"):
        process_image_file(filename, filepath, ext)
        return

    if ext in AUDIO_VIDEO_TYPES:
        print(f"  Audio/video — transcription required.")
        file_hash = get_file_hash(filepath) or "—"
        with open(INBOX_MD, "a") as f:
            f.write(f"| {inbox_safe(filename)} | {datetime.date.today()} | {file_hash} | — | TRANSCRIPTION NEEDED |\n")
        move_to_processed(filename)
        return

    if ext in LEGACY_OFFICE_TYPES:
        content = extract_legacy_office(filepath, filename, ext)
        if not content:
            log_unprocessable(
                filename,
                f"legacy .{ext} format — re-save as .{LEGACY_TO_MODERN[ext]} "
                "or export to PDF and re-drop (LibreOffice auto-conversion unavailable or failed)",
            )
            return
    else:
        content = read_inbox_file(filename)
        if not content:
            log_unprocessable(filename, "no text could be extracted (see run log for details)")
            return

    print(f"  Running ingest agent...")
    result = run_ingest_agent(filename, content)
    if not result:
        # Do NOT move the file: leaving it in the inbox means the next push or
        # the weekly sweep retries it. Moving it with no inbox record would
        # silently lose it forever.
        log_unprocessable(filename, "ingest agent failed after retries (see run log)",
                          permanent=False)
        return
    apply_vault_changes(result, filename)
    move_to_processed(filename)
    print(f"  ✓ Done: {filename}")

def main():
    print("\n" + "=" * 60)
    print(f"ABAP Vault Ingest — {datetime.date.today()}")
    print("=" * 60 + "\n")

    git_pull()

    payload = get_payload()

    if payload:
        process_file(payload["filename"])
    else:
        print("Safety-net scan: checking raw/inbox/...")
        unprocessed = scan_inbox_for_unprocessed()
        if unprocessed:
            print(f"Found {len(unprocessed)} file(s): {unprocessed}")
            for filename in unprocessed:
                process_file(filename)
        else:
            print("Nothing to process.")

    print("\nDone.")

if __name__ == "__main__":
    main()
