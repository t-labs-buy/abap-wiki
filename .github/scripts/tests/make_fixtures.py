#!/usr/bin/env python3
"""Build hard-case fixtures for the extractor A/B test.

Every fixture embeds sentinel strings whose presence/position the harness
asserts on. Hard cases covered:
  - empty table cells (column alignment must survive)
  - merged cells (horizontal + vertical)
  - document order (heading -> para -> table -> para interleaving)
  - speaker notes + grouped shapes (PPTX)
  - uncached formulas (XLSX written by openpyxl, never opened in Excel)
  - scanned PDF (image-only, no text layer)
  - legacy .xls
"""
import os

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "fixtures")
os.makedirs(OUT, exist_ok=True)


def make_docx():
    from docx import Document
    doc = Document()
    doc.add_heading("BP Address Validation", level=1)
    doc.add_paragraph("INTRO-BEFORE-TABLE: scope covers QAS and PRD.")
    t = doc.add_table(rows=4, cols=3)
    # header
    t.cell(0, 0).text = "Field"
    t.cell(0, 1).text = "Mandatory"
    t.cell(0, 2).text = "Default"
    # row with EMPTY MIDDLE CELL: Street | <empty> | NONE_DEFAULT
    t.cell(1, 0).text = "Street"
    t.cell(1, 2).text = "NONE_DEFAULT"
    # horizontally merged row: "MERGED-SPAN-NOTE" spans col 0-2
    m = t.cell(2, 0).merge(t.cell(2, 1)).merge(t.cell(2, 2))
    m.text = "MERGED-SPAN-NOTE applies to all fields"
    # normal row
    t.cell(3, 0).text = "City"
    t.cell(3, 1).text = "Yes"
    t.cell(3, 2).text = "BERLIN_DEFAULT"
    doc.add_paragraph("OUTRO-AFTER-TABLE: validation runs in BAdI.")
    doc.add_heading("Second Section", level=2)
    doc.add_paragraph("TAIL-PARA content.")
    doc.save(os.path.join(OUT, "hard.docx"))


def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    # Slide 1: title + table with empty + merged cells, plus speaker notes
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    s1.shapes.title.text = "SLIDE-ONE-TITLE Wave 2 Estimates"
    rows, cols = 3, 3
    tbl = s1.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(2)).table
    tbl.cell(0, 0).text = "Object"
    tbl.cell(0, 1).text = "Days"
    tbl.cell(0, 2).text = "Owner"
    tbl.cell(1, 0).text = "ZSD_ORDER_CHECK"
    # (1,1) intentionally EMPTY
    tbl.cell(1, 2).text = "ANNA_OWNER"
    tbl.cell(2, 0).merge(tbl.cell(2, 2))
    tbl.cell(2, 0).text = "MERGED-FOOTER-ROW totals pending"
    s1.notes_slide.notes_text_frame.text = "SPEAKER-NOTE-ALPHA: client worried about VA01 performance."
    # Slide 2: grouped shapes
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.util import Pt
    box1 = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    box1.text_frame.text = "GROUP-CHILD-ONE"
    box2 = s2.shapes.add_textbox(Inches(1), Inches(2.5), Inches(3), Inches(1))
    box2.text_frame.text = "GROUP-CHILD-TWO"
    g = s2.shapes.add_group_shape([box1, box2])
    s2.notes_slide.notes_text_frame.text = "SPEAKER-NOTE-BETA on slide two."
    prs.save(os.path.join(OUT, "hard.pptx"))


def make_xlsx():
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Estimates"
    ws.append(["Object", "Est", "Act", "Var"])
    ws.append(["ZSD_ORDER_CHECK", 5, 7, "=C2-B2"])   # formula, NO cached value
    ws.append(["ZMM_GR_BADI", None, 3, ""])           # empty mid-row cell
    ws["A5"] = "TOTAL_ROW"
    ws["B5"] = "=SUM(B2:B3)"                          # uncached formula
    ws.merge_cells("A7:C7")
    ws["A7"] = "MERGED-XLSX-NOTE assumptions apply"
    ws2 = wb.create_sheet("Assumptions")
    ws2.append(["SECOND-SHEET-MARKER", "unit test env available"])
    wb.save(os.path.join(OUT, "hard.xlsx"))


def make_xls():
    import xlwt
    wb = xlwt.Workbook()
    ws = wb.add_sheet("Legacy")
    ws.write(0, 0, "LEGACY-XLS-HEADER")
    ws.write(0, 1, "Days")
    ws.write(1, 0, "ZFI_OLD_REPORT")
    # (1,1) empty
    ws.write(1, 2, "LEGACY-TAIL-CELL")
    wb.save(os.path.join(OUT, "legacy.xls"))


def make_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(os.path.join(OUT, "hard.pdf"), pagesize=A4)
    data = [
        ["Field", "Mandatory", "Default"],
        ["Street", "", "PDF_NONE_DEFAULT"],      # empty middle cell
        ["City", "Yes", "PDF_BERLIN"],
    ]
    t = Table(data)
    t.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 0.5, colors.black)]))
    doc.build([
        Paragraph("PDF-BEFORE-TABLE heading text", styles["Heading1"]),
        Spacer(1, 12), t, Spacer(1, 12),
        Paragraph("PDF-AFTER-TABLE trailing paragraph", styles["Normal"]),
    ])


def make_scanned_pdf():
    """Image-only PDF: render text into a PIL image, embed as full page."""
    from PIL import Image, ImageDraw, ImageFont
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    img_path = os.path.join(OUT, "_scan_page.png")
    img = Image.new("RGB", (1654, 2339), "white")
    d = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 48)
    except Exception:
        font = ImageFont.load_default()
    d.text((120, 200), "SCANNED-OCR-MARKER Transport Release Runbook", fill="black", font=font)
    d.text((120, 320), "Step 1: release task in SE09", fill="black", font=font)
    d.text((120, 440), "Step 2: import into QAS via STMS", fill="black", font=font)
    img.save(img_path)
    c = canvas.Canvas(os.path.join(OUT, "scanned.pdf"), pagesize=A4)
    c.drawImage(img_path, 0, 0, width=A4[0], height=A4[1])
    c.showPage()
    c.save()
    os.remove(img_path)


if __name__ == "__main__":
    make_docx(); make_pptx(); make_xlsx(); make_xls(); make_pdf(); make_scanned_pdf()
    for f in sorted(os.listdir(OUT)):
        print(f, os.path.getsize(os.path.join(OUT, f)))
